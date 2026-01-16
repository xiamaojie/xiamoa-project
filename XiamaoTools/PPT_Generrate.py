import streamlit as st
import os
import random
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image, ImageDraw, ImageFont
import datetime
import sys

# 确保输出目录存在
OUTPUT_DIR = "/Users/admin/TestLog/ppt_data"
if not os.path.exists(OUTPUT_DIR):
    os.makedirs(OUTPUT_DIR)
    print(f"创建目录: {OUTPUT_DIR}")


# 生成随机文本内容
def generate_random_text(length=100):
    words = ["PPT", "测试", "数据", "生成", "工具", "幻灯片", "内容", "标题", "项目", "要点",
             "演示", "设计", "自动", "创建", "大型", "文件", "示例", "文本", "随机", "填充",
             "人工智能", "机器学习", "数据分析", "可视化", "企业", "解决方案", "技术", "创新",
             "发展", "市场", "产品", "服务", "客户", "价值", "战略", "团队", "协作", "效率", "提升"]
    return ' '.join(random.choices(words, k=length))


# 生成随机图表（作为图片）
def generate_random_chart():
    fig, ax = plt.subplots(figsize=(6, 4))
    categories = ['Q1', 'Q2', 'Q3', 'Q4']
    values = [random.randint(50, 100) for _ in range(4)]
    ax.bar(categories, values, color=plt.cm.viridis(np.linspace(0, 1, 4)))
    ax.set_title('季度业绩报告', fontsize=10)
    ax.set_ylabel('销售额 (万元)')
    fig.tight_layout()

    temp_file = os.path.join(OUTPUT_DIR, f"temp_chart_{time.time()}.png")
    plt.savefig(temp_file, dpi=100)
    plt.close()
    return temp_file


# 生成随机图像
def generate_random_image():
    img = Image.new('RGB', (800, 600), color=(random.randint(200, 255),
                                              random.randint(200, 255),
                                              random.randint(200, 255)))
    draw = ImageDraw.Draw(img)

    # 添加一些随机形状
    for _ in range(5):
        x1, y1 = random.randint(0, 700), random.randint(0, 500)
        x2, y2 = x1 + random.randint(50, 200), y1 + random.randint(50, 200)
        draw.rectangle([x1, y1, x2, y2],
                       fill=(random.randint(50, 200),
                             random.randint(50, 200),
                             random.randint(50, 200)),
                       outline=(0, 0, 0))

    # 添加一些文字
    try:
        font = ImageFont.truetype("arial.ttf", 24)
    except:
        font = ImageFont.load_default()

    for i in range(3):
        text = f"图像 #{i + 1}"
        x = random.randint(50, 600)
        y = random.randint(50, 500)
        draw.text((x, y), text, fill=(0, 0, 0), font=font)

    temp_file = os.path.join(OUTPUT_DIR, f"temp_img_{time.time()}.png")
    img.save(temp_file)
    return temp_file


# 创建大型PPT文件
def create_large_ppt(num_slides=50, include_images=True, include_charts=True, include_tables=True):
    start_time = time.time()

    # 创建演示文稿
    prs = Presentation()

    # 进度条
    progress_bar = st.progress(0)
    status_text = st.empty()

    # 定义布局
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    section_header_layout = prs.slide_layouts[2]
    two_content_layout = prs.slide_layouts[3]

    # 创建幻灯片
    for i in range(num_slides):
        # 更新进度
        percent_complete = (i + 1) / num_slides
        progress_bar.progress(percent_complete)
        status_text.text(f"正在生成幻灯片 {i + 1}/{num_slides}...")

        # 每10张幻灯片使用一个不同的布局
        if i % 10 == 0:
            slide_layout = section_header_layout
        elif i % 2 == 0:
            slide_layout = two_content_layout
        else:
            slide_layout = content_slide_layout

        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = slide.shapes.title
        title.text = f"幻灯片 {i + 1} - {generate_random_text(5)}"
        title.text_frame.paragraphs[0].font.size = Pt(32 if i % 10 == 0 else 24)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(
            random.randint(0, 100),
            random.randint(0, 100),
            random.randint(100, 150)
        )

        # 添加内容
        if slide_layout == content_slide_layout:
            content = slide.placeholders[1]
            content.text = generate_random_text(100)
            for paragraph in content.text_frame.paragraphs:
                paragraph.font.size = Pt(18)

        elif slide_layout == two_content_layout:
            left_content = slide.placeholders[1]
            left_content.text = generate_random_text(50)
            for paragraph in left_content.text_frame.paragraphs:
                paragraph.font.size = Pt(16)

            right_content = slide.placeholders[2]
            right_content.text = generate_random_text(50)
            for paragraph in right_content.text_frame.paragraphs:
                paragraph.font.size = Pt(16)

        # 添加图像
        if include_images and i % 3 == 0:
            img_path = generate_random_image()
            left = Inches(random.uniform(0.5, 5))
            top = Inches(random.uniform(1.5, 4))
            slide.shapes.add_picture(img_path, left, top, width=Inches(3))
            os.unlink(img_path)  # 删除临时图像文件

        # 添加图表
        if include_charts and i % 4 == 0:
            chart_path = generate_random_chart()
            left = Inches(random.uniform(0.5, 5))
            top = Inches(random.uniform(1.5, 4))
            slide.shapes.add_picture(chart_path, left, top, width=Inches(4))
            os.unlink(chart_path)  # 删除临时图表文件

        # 添加表格
        if include_tables and i % 5 == 0:
            rows = random.randint(4, 8)
            cols = random.randint(3, 6)
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(9)
            height = Inches(0.5 * rows)

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            # 设置列宽
            for col_idx in range(cols):
                table.columns[col_idx].width = Inches(width.inches / cols)

            # 填充表格
            for r in range(rows):
                for c in range(cols):
                    cell = table.cell(r, c)
                    if r == 0:  # 表头
                        cell.text = f"列 {c + 1}"
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(100, 150, 200)
                        cell.text_frame.paragraphs[0].font.bold = True
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    else:
                        cell.text = str(random.randint(1, 1000))
                        cell.text_frame.paragraphs[0].font.size = Pt(12)

    # 保存PPT到指定目录
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"large_presentation_{num_slides}slides_{timestamp}.pptx"
    file_path = os.path.join(OUTPUT_DIR, filename)
    prs.save(file_path)

    # 计算文件大小
    file_size = os.path.getsize(file_path) / (1024 * 1024)  # 转换为MB

    end_time = time.time()
    generation_time = end_time - start_time

    status_text.text(f"生成完成！耗时: {generation_time:.2f}秒, 文件大小: {file_size:.2f} MB")
    progress_bar.empty()

    st.success(f"PPT文件已保存到: {file_path}")

    return file_path


# Streamlit界面
def main():
    st.title("📊 大型PPT测试文件生成工具")
    st.write(f"此工具可生成包含大量幻灯片、文本、图表和图像的PPT文件，文件将保存到: `{OUTPUT_DIR}`")

    with st.expander("⚙️ 生成选项", expanded=True):
        col1, col2 = st.columns(2)
        with col1:
            num_slides = st.slider("幻灯片数量", 10, 500, 100)
            include_images = st.checkbox("包含随机图像", value=True)
        with col2:
            include_charts = st.checkbox("包含随机图表", value=True)
            include_tables = st.checkbox("包含随机表格", value=True)

    if st.button("生成PPT文件", type="primary", use_container_width=True):
        with st.spinner("正在创建大型PPT文件，请稍候..."):
            ppt_file = create_large_ppt(
                num_slides=num_slides,
                include_images=include_images,
                include_charts=include_charts,
                include_tables=include_tables
            )

            # 提供下载链接
            with open(ppt_file, "rb") as f:
                ppt_bytes = f.read()

            st.download_button(
                label="下载PPT文件",
                data=ppt_bytes,
                file_name=os.path.basename(ppt_file),
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )

    # 显示目录中的文件列表
    st.markdown("---")
    st.subheader("已生成的文件列表")
    if os.path.exists(OUTPUT_DIR):
        files = os.listdir(OUTPUT_DIR)
        ppt_files = [f for f in files if f.endswith('.pptx')]

        if ppt_files:
            st.write(f"目录 `{OUTPUT_DIR}` 中的PPT文件:")
            for file in ppt_files:
                file_path = os.path.join(OUTPUT_DIR, file)
                file_size = os.path.getsize(file_path) / (1024 * 1024)
                modified_time = datetime.datetime.fromtimestamp(os.path.getmtime(file_path)).strftime(
                    '%Y-%m-%d %H:%M:%S')

                col1, col2, col3 = st.columns([6, 2, 2])
                col1.write(file)
                col2.write(f"{file_size:.2f} MB")
                col3.write(modified_time)

                # 提供删除按钮
                if st.button(f"删除 {file}", key=f"del_{file}"):
                    try:
                        os.remove(file_path)
                        st.success(f"已删除: {file}")
                        st.experimental_rerun()
                    except Exception as e:
                        st.error(f"删除失败: {str(e)}")
        else:
            st.info("该目录中还没有PPT文件")
    else:
        st.error(f"目录不存在: {OUTPUT_DIR}")


# 添加直接运行支持
if __name__ == '__main__':
    if 'streamlit' in sys.modules:
        main()
    else:
        import subprocess

        subprocess.run(["streamlit", "run", sys.argv[0]])